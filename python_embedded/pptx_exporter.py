"""
PowerPoint Export Module for easyCris Plots

Exports plots to PowerPoint format with embedded data in slide notes
for later recovery. Uses python-pptx library.

Version: 1.0.0
"""

import json
import sys
import io
import base64
from pathlib import Path

# Add python_dependencies to path
script_dir = Path(__file__).parent
deps_dir = script_dir / "python_dependencies"
if str(deps_dir) not in sys.path:
    sys.path.insert(0, str(deps_dir))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN


def export_plot_to_pptx(
    image_base64: str,
    plot_data: dict,
    output_path: str,
    options: dict = None
) -> dict:
    """
    Export plot to PowerPoint with embedded data in notes.

    Args:
        image_base64: Base64-encoded PNG image of the plot
        plot_data: PlotSpec dictionary containing all plot data
        output_path: Path to save the .pptx file
        options: Export options (embed_data, add_branding, slide_layout)

    Returns:
        dict with success status and message
    """
    options = options or {}
    embed_data = options.get("embed_data", True)
    add_branding = options.get("add_branding", True)

    try:
        # Decode image
        image_bytes = base64.b64decode(image_base64)

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)

        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add plot image (centered, with margins)
        image_stream = io.BytesIO(image_bytes)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.333)  # Leave 0.5" margins

        slide.shapes.add_picture(image_stream, left, top, width=width)

        # Add title if available
        title = plot_data.get("title", "easyCris Plot")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER

        # Add branding footer
        if add_branding:
            footer_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3)
            )
            footer_frame = footer_box.text_frame
            footer_para = footer_frame.paragraphs[0]
            footer_para.text = "Created with easyCris | Data recoverable from slide notes"
            footer_para.font.size = Pt(10)
            footer_para.font.italic = True
            footer_para.alignment = PP_ALIGN.CENTER

        # Embed data in slide notes for recovery
        if embed_data:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame

            # Serialize plot data as JSON
            data_json = json.dumps(plot_data, indent=2, default=str)
            notes_frame.text = f"easyCris_DATA:{data_json}"

        # Set document properties
        prs.core_properties.title = title
        prs.core_properties.subject = f"easyCris Plot - {plot_data.get('type', 'unknown')}"
        prs.core_properties.keywords = "easyCris,recoverable,statistical-plot"
        prs.core_properties.comments = "This presentation contains recoverable easyCris plot data in slide notes."

        # Save presentation
        prs.save(output_path)

        return {
            "success": True,
            "message": f"Plot exported to {output_path}",
            "path": output_path
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "message": f"Failed to export plot: {str(e)}"
        }


def recover_data_from_pptx(pptx_path: str) -> dict:
    """
    Extract easyCris plot data from a PowerPoint file.

    Args:
        pptx_path: Path to the .pptx file

    Returns:
        dict with success status and recovered plots
    """
    try:
        prs = Presentation(pptx_path)
        recovered_plots = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            try:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text

                if notes_text.startswith("easyCris_DATA:"):
                    json_str = notes_text[len("easyCris_DATA:"):]
                    plot_data = json.loads(json_str)
                    recovered_plots.append({
                        "slide": slide_num,
                        "data": plot_data,
                        "title": plot_data.get("title", f"Slide {slide_num}")
                    })
            except Exception:
                # Skip slides without valid notes
                continue

        return {
            "success": True,
            "plots": recovered_plots,
            "count": len(recovered_plots),
            "message": f"Recovered {len(recovered_plots)} plot(s) from {pptx_path}"
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "plots": [],
            "count": 0,
            "message": f"Failed to recover data: {str(e)}"
        }


def export_multiple_plots_to_pptx(
    plots: list,
    output_path: str,
    options: dict = None
) -> dict:
    """
    Export multiple plots to a single PowerPoint presentation.

    Args:
        plots: List of dicts, each with 'image_base64' and 'plot_data'
        output_path: Path to save the .pptx file
        options: Export options

    Returns:
        dict with success status and message
    """
    options = options or {}
    embed_data = options.get("embed_data", True)
    add_branding = options.get("add_branding", True)

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        for plot in plots:
            image_base64 = plot.get("image_base64", "")
            plot_data = plot.get("plot_data", {})

            if not image_base64:
                continue

            # Decode image
            image_bytes = base64.b64decode(image_base64)
            slide = prs.slides.add_slide(blank_layout)

            # Add image
            image_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                image_stream, Inches(0.5), Inches(0.5), width=Inches(12.333)
            )

            # Add title
            title = plot_data.get("title", "easyCris Plot")
            if title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.4)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = title
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER

            # Add branding
            if add_branding:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3)
                )
                footer_frame = footer_box.text_frame
                footer_para = footer_frame.paragraphs[0]
                footer_para.text = "Created with easyCris | Data recoverable from slide notes"
                footer_para.font.size = Pt(10)
                footer_para.font.italic = True
                footer_para.alignment = PP_ALIGN.CENTER

            # Embed data in notes
            if embed_data:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                data_json = json.dumps(plot_data, indent=2, default=str)
                notes_frame.text = f"easyCris_DATA:{data_json}"

        # Set document properties
        prs.core_properties.title = "easyCris Plots"
        prs.core_properties.subject = f"Collection of {len(plots)} statistical plots"
        prs.core_properties.keywords = "easyCris,recoverable,statistical-plot"

        prs.save(output_path)

        return {
            "success": True,
            "message": f"Exported {len(plots)} plots to {output_path}",
            "path": output_path,
            "count": len(plots)
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "message": f"Failed to export plots: {str(e)}"
        }


# CLI interface for Tauri invocation
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="easyCris PPTX Exporter")
    parser.add_argument("--action", required=True, choices=["export", "recover", "export_multi"])
    parser.add_argument("--input", help="Input file or JSON data")
    parser.add_argument("--output", help="Output file path")
    parser.add_argument("--embed-data", action="store_true", default=True)
    parser.add_argument("--no-embed-data", action="store_false", dest="embed_data")
    parser.add_argument("--branding", action="store_true", default=True)
    parser.add_argument("--no-branding", action="store_false", dest="branding")

    args = parser.parse_args()

    if args.action == "export":
        # Read input JSON from stdin
        input_data = json.loads(sys.stdin.read())
        result = export_plot_to_pptx(
            image_base64=input_data.get("image_base64", ""),
            plot_data=input_data.get("plot_data", {}),
            output_path=args.output,
            options={
                "embed_data": args.embed_data,
                "add_branding": args.branding
            }
        )
        print(json.dumps(result))

    elif args.action == "recover":
        result = recover_data_from_pptx(args.input)
        print(json.dumps(result))

    elif args.action == "export_multi":
        input_data = json.loads(sys.stdin.read())
        result = export_multiple_plots_to_pptx(
            plots=input_data.get("plots", []),
            output_path=args.output,
            options={
                "embed_data": args.embed_data,
                "add_branding": args.branding
            }
        )
        print(json.dumps(result))
